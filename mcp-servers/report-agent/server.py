"""
Report Agent MCP Server
Handles PDF and PowerPoint report generation
"""
import asyncio
import json
from pathlib import Path
from typing import Any, Dict
from datetime import datetime
import html

from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib import colors

from pptx import Presentation
from pptx.util import Inches, Pt


class ReportAgent:
    """Agent for report generation tasks"""
    
    def __init__(self):
        self.output_dir = Path("./reports")
    
    @staticmethod
    def escape_html(text: str) -> str:
        """Escape HTML special characters for ReportLab Paragraph"""
        if not text:
            return ""
        # Escape HTML entities to prevent ReportLab parsing errors
        return html.escape(str(text))
    
    async def create_pdf_report(self, content: Dict[str, Any], output_path: str) -> Dict[str, Any]:
        """
        Create a PDF report from analysis results
        
        Args:
            content: Report content with sections
            output_path: Path for output PDF file
            
        Returns:
            Result with PDF file path
        """
        try:
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Create PDF
            doc = SimpleDocTemplate(
                str(output_file),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )
            
            # Container for elements
            elements = []
            styles = getSampleStyleSheet()
            
            # Title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2C3E50'),
                spaceAfter=12,
                alignment=1  # Center
            )
            
            report_title = content.get("title", "Video Analysis Report")
            elements.append(Paragraph(self.escape_html(report_title), title_style))
            
            # Subtitle with generated timestamp
            subtitle_style = ParagraphStyle(
                'Subtitle',
                parent=styles['Normal'],
                fontSize=12,
                textColor=colors.HexColor('#7F8C8D'),
                spaceAfter=30,
                alignment=1  # Center
            )
            generated_at = content.get("generated_at", datetime.now().strftime("%Y-%m-%d %I:%M %p"))
            elements.append(Paragraph(f"Generated: {self.escape_html(generated_at)}", subtitle_style))
            elements.append(Spacer(1, 0.3 * inch))
            
            # Define section styles
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=16,
                textColor=colors.HexColor('#2C3E50'),
                spaceAfter=12,
                spaceBefore=12
            )
            
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=6
            )
            
            # SUMMARY Section
            summary = content.get("summary", "")
            if summary:
                elements.append(Paragraph("SUMMARY:", heading2_style))
                elements.append(Paragraph(self.escape_html(summary), normal_style))
                elements.append(Spacer(1, 0.2 * inch))
            
            # KEY MOMENTS Section
            key_moments = content.get("key_moments", [])
            if key_moments:
                elements.append(Paragraph("KEY MOMENTS:", heading2_style))
                for moment in key_moments:
                    if isinstance(moment, dict):
                        time = self.escape_html(moment.get("time", ""))
                        event = self.escape_html(moment.get("event", ""))
                        elements.append(Paragraph(f"{time} - {event}", normal_style))
                    else:
                        elements.append(Paragraph(self.escape_html(str(moment)), normal_style))
                elements.append(Spacer(1, 0.2 * inch))
            
            # VISUAL ANALYSIS Section
            visual_analysis = content.get("visual_analysis", "")
            if visual_analysis:
                elements.append(Paragraph("VISUAL ANALYSIS:", heading2_style))
                # Split by newlines and format each line
                for line in visual_analysis.split('\n'):
                    if line.strip():
                        elements.append(Paragraph(self.escape_html(line), normal_style))
                elements.append(Spacer(1, 0.2 * inch))
            
            # TRANSCRIPT Section
            transcript = content.get("transcript", "")
            if transcript:
                elements.append(Paragraph("TRANSCRIPT:", heading2_style))
                elements.append(Paragraph(self.escape_html(transcript), normal_style))
                elements.append(Spacer(1, 0.2 * inch))
            
            # KEY TAKEAWAYS Section
            takeaways = content.get("takeaways", [])
            if takeaways:
                elements.append(Paragraph("KEY TAKEAWAYS:", heading2_style))
                for takeaway in takeaways:
                    elements.append(Paragraph(f"• {self.escape_html(takeaway)}", normal_style))
                elements.append(Spacer(1, 0.2 * inch))
            
            # Build PDF
            doc.build(elements)
            
            return {
                "success": True,
                "output_path": str(output_file),
                "size_bytes": output_file.stat().st_size
            }
            
        except Exception as e:
            return {"error": f"PDF creation failed: {str(e)}"}
    
    async def create_ppt_report(self, content: Dict[str, Any], output_path: str) -> Dict[str, Any]:
        """
        Create a PowerPoint report from analysis results
        
        Args:
            content: Report content with sections
            output_path: Path for output PPTX file
            
        Returns:
            Result with PPTX file path
        """
        try:
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            report_title = content.get("title", "Video Analysis Report")
            title.text = report_title
            generated_at = content.get("generated_at", datetime.now().strftime("%Y-%m-%d %I:%M %p"))
            subtitle.text = f"Generated: {generated_at}"
            
            # SUMMARY Slide
            summary = content.get("summary", "")
            if summary:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "SUMMARY"
                
                tf = body_shape.text_frame
                tf.text = summary
            
            # KEY MOMENTS Slide
            key_moments = content.get("key_moments", [])
            if key_moments:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "KEY MOMENTS"
                
                tf = body_shape.text_frame
                first = True
                for moment in key_moments:
                    if isinstance(moment, dict):
                        time = moment.get("time", "")
                        event = moment.get("event", "")
                        text = f"{time} - {event}"
                    else:
                        text = str(moment)
                    
                    if first:
                        tf.text = text
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = text
                        p.level = 0
            
            # VISUAL ANALYSIS Slide
            visual_analysis = content.get("visual_analysis", "")
            if visual_analysis:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "VISUAL ANALYSIS"
                
                tf = body_shape.text_frame
                # Use first 500 chars to avoid overcrowding
                tf.text = visual_analysis[:500]
            
            # TRANSCRIPT Slide
            transcript = content.get("transcript", "")
            if transcript:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "TRANSCRIPT"
                
                tf = body_shape.text_frame
                # Use first 500 chars to avoid overcrowding
                tf.text = transcript[:500]
            
            # KEY TAKEAWAYS Slide
            takeaways = content.get("takeaways", [])
            if takeaways:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = "KEY TAKEAWAYS"
                
                tf = body_shape.text_frame
                first = True
                for takeaway in takeaways:
                    if first:
                        tf.text = f"• {takeaway}"
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = f"• {takeaway}"
                        p.level = 0
            
            # Save presentation
            prs.save(str(output_file))
            
            return {
                "success": True,
                "output_path": str(output_file),
                "size_bytes": output_file.stat().st_size
            }
            
        except Exception as e:
            return {"error": f"PowerPoint creation failed: {str(e)}"}
    
    async def format_content(self, raw_content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Format and structure content for reports
        
        Args:
            raw_content: Raw analysis results
            
        Returns:
            Formatted content ready for report generation
        """
        try:
            formatted = {
                "video_id": raw_content.get("video_id", "unknown"),
                "metadata": raw_content.get("metadata", {}),
                "sections": ["summary", "transcription", "key_frames"],
                "summary": "Video analysis summary would go here.",
                "transcription": raw_content.get("transcription", {}).get("text", "No transcription available"),
                "key_frames": f"{len(raw_content.get('frames', []))} key frames extracted"
            }
            
            return {
                "success": True,
                "formatted_content": formatted
            }
            
        except Exception as e:
            return {"error": f"Content formatting failed: {str(e)}"}


async def main():
    """Main entry point for the MCP server"""
    agent = ReportAgent()
    server = Server("report-agent")
    
    # Register tools
    @server.list_tools()
    async def list_tools() -> list[Tool]:
        return [
            Tool(
                name="create_pdf_report",
                description="Generate a PDF report from video analysis results",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "content": {
                            "type": "object",
                            "description": "Report content with metadata and sections"
                        },
                        "output_path": {
                            "type": "string",
                            "description": "Path for output PDF file"
                        }
                    },
                    "required": ["content", "output_path"]
                }
            ),
            Tool(
                name="create_ppt_report",
                description="Generate a PowerPoint presentation from video analysis results",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "content": {
                            "type": "object",
                            "description": "Report content with metadata and sections"
                        },
                        "output_path": {
                            "type": "string",
                            "description": "Path for output PPTX file"
                        }
                    },
                    "required": ["content", "output_path"]
                }
            ),
            Tool(
                name="format_content",
                description="Format and structure raw analysis results for reports",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "raw_content": {
                            "type": "object",
                            "description": "Raw analysis results to format"
                        }
                    },
                    "required": ["raw_content"]
                }
            )
        ]
    
    @server.call_tool()
    async def call_tool(name: str, arguments: Any) -> list[TextContent]:
        if name == "create_pdf_report":
            result = await agent.create_pdf_report(
                arguments["content"],
                arguments["output_path"]
            )
        elif name == "create_ppt_report":
            result = await agent.create_ppt_report(
                arguments["content"],
                arguments["output_path"]
            )
        elif name == "format_content":
            result = await agent.format_content(arguments["raw_content"])
        else:
            result = {"error": f"Unknown tool: {name}"}
        
        return [TextContent(type="text", text=json.dumps(result, indent=2))]
    
    # Run server
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


if __name__ == "__main__":
    asyncio.run(main())
